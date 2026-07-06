# -*- coding: utf-8 -*-
# =============================================================================
# Ag-ppt-create - AI-powered PPTX generation pipeline
# https://github.com/aktsmm/Ag-ppt-create
# 
# Copyright (c) aktsmm. Licensed under CC BY-NC-SA 4.0.
# DO NOT MODIFY THIS HEADER BLOCK.
# =============================================================================
"""
Clean PPTX template by removing problematic elements.

Usage:
    python scripts/clean_template.py <input.pptx> <output.pptx> [--keep-first-master] [--keep-metadata]

Removes / sanitizes:
    - Background images (PICTURE shapes) from slide masters
    - Blip references from Picture Placeholders in layouts
    - docProps/custom.xml (MIP labels, tenant IDs, SharePoint columns,
      colleague identities) and cp:lastModifiedBy / cp:revision in
      docProps/core.xml, so the check-in artifact does not carry the
      metadata that PowerPoint attaches when the file is opened from
      OneDrive/SharePoint. Use --keep-metadata only for tenant-internal
      private repos where the labels are intentional.
    - Optionally keeps only the first slide master

This prepares external templates (e.g., from Microsoft Ignite) for use
without decorative background images interfering with generated content.
"""

from pptx import Presentation
from lxml import etree
import sys
import zipfile
from pathlib import Path
import shutil


def remove_master_background_images(prs) -> int:
    """Remove PICTURE shapes from slide masters."""
    removed = 0
    for master in prs.slide_masters:
        shapes_to_remove = []
        for shape in master.shapes:
            if 'PICTURE' in str(shape.shape_type):
                shapes_to_remove.append(shape)
        
        for shape in shapes_to_remove:
            try:
                sp = shape._element
                sp.getparent().remove(sp)
                removed += 1
            except Exception as e:
                print(f"  [!] Could not remove shape: {e}")
    
    return removed


def remove_placeholder_blips(prs) -> int:
    """Remove blip elements from Picture Placeholders in layouts."""
    removed = 0
    ns = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    }
    
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            for shape in layout.shapes:
                if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                    try:
                        ph_type = str(shape.placeholder_format.type)
                        if 'PICTURE' in ph_type:
                            # Find and remove blip elements
                            pic = shape._element
                            blip_elements = pic.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}blip')
                            for blip in blip_elements:
                                blip.getparent().remove(blip)
                                removed += 1
                    except (ValueError, AttributeError):
                        pass
    
    return removed


def sanitize_pptx_metadata(pptx_path: str) -> dict:
    """Strip PII from docProps in a saved pptx (in-place).

    - docProps/custom.xml -> empty <Properties/> element (removes MIP /
      MSIP labels, SharePoint columns, tenant IDs, colleague emails).
    - docProps/core.xml   -> blank cp:lastModifiedBy and cp:revision.

    PowerPoint re-instruments docProps on every save from OneDrive /
    SharePoint by a labeled account, so this must run on the check-in
    artifact whenever the template is regenerated, not once per skill.
    """
    EMPTY_CUSTOM = (
        b'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n'
        b'<Properties xmlns="http://schemas.openxmlformats.org/officeDocument/2006/custom-properties" '
        b'xmlns:vt="http://schemas.openxmlformats.org/officeDocument/2006/docPropsVTypes"/>'
    )
    CORE_NS = {
        'cp': 'http://schemas.openxmlformats.org/package/2006/metadata/core-properties',
    }

    stats = {'custom_stripped': False, 'core_neutralized': False}
    src = Path(pptx_path)
    tmp = src.with_suffix(src.suffix + '.sanitize-tmp')

    try:
        with zipfile.ZipFile(src, 'r') as zin, zipfile.ZipFile(
            tmp, 'w', compression=zipfile.ZIP_DEFLATED
        ) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename == 'docProps/custom.xml':
                    if data.strip() and data != EMPTY_CUSTOM:
                        stats['custom_stripped'] = True
                    data = EMPTY_CUSTOM
                elif item.filename == 'docProps/core.xml':
                    try:
                        root = etree.fromstring(data)
                        changed = False
                        for xp in ('cp:lastModifiedBy', 'cp:revision'):
                            el = root.find(xp, CORE_NS)
                            if el is not None and (el.text or '').strip():
                                el.text = ''
                                changed = True
                        if changed:
                            stats['core_neutralized'] = True
                            data = etree.tostring(
                                root, xml_declaration=True,
                                encoding='UTF-8', standalone=True
                            )
                    except etree.XMLSyntaxError:
                        pass
                zout.writestr(item, data)
        shutil.move(str(tmp), str(src))
    finally:
        if tmp.exists():
            try:
                tmp.unlink()
            except OSError:
                pass

    return stats


def clean_template(input_path: str, output_path: str, keep_first_master_only: bool = False, sanitize_metadata: bool = True) -> dict:
    """
    Clean PPTX template and save to output path.
    
    Args:
        input_path: Path to input PPTX
        output_path: Path to output PPTX
        keep_first_master_only: If True, remove all but the first slide master
        sanitize_metadata: If True (default), strip docProps PII (MIP labels,
            SharePoint fields, cp:lastModifiedBy, cp:revision) after save.
    
    Returns:
        dict with cleaning statistics
    """
    prs = Presentation(input_path)
    
    stats = {
        'input': Path(input_path).name,
        'output': Path(output_path).name,
        'masters_before': len(prs.slide_masters),
        'background_images_removed': 0,
        'placeholder_blips_removed': 0,
        'masters_removed': 0
    }
    
    print(f"\n🔧 Cleaning template: {input_path}")
    print(f"   Masters: {stats['masters_before']}")
    
    # Step 1: Remove background images from masters
    bg_removed = remove_master_background_images(prs)
    stats['background_images_removed'] = bg_removed
    print(f"   Removed {bg_removed} background image(s) from masters")
    
    # Step 2: Remove blip references from placeholders
    blip_removed = remove_placeholder_blips(prs)
    stats['placeholder_blips_removed'] = blip_removed
    print(f"   Removed {blip_removed} blip reference(s) from placeholders")
    
    # Step 3: Optionally remove extra masters (not fully implemented - complex)
    # This would require careful handling of slide relationships
    if keep_first_master_only and len(prs.slide_masters) > 1:
        print(f"   [i] Note: --keep-first-master not fully implemented yet")
    
    # Save cleaned template
    prs.save(output_path)
    
    total_removed = bg_removed + blip_removed
    stats['total_removed'] = total_removed
    
    print(f"\n✅ Saved clean template: {output_path}")
    print(f"   Total elements removed: {total_removed}")

    # Strip PII from docProps unless the caller opted out
    if sanitize_metadata:
        meta = sanitize_pptx_metadata(output_path)
        stats['metadata_custom_stripped'] = meta['custom_stripped']
        stats['metadata_core_neutralized'] = meta['core_neutralized']
        if meta['custom_stripped'] or meta['core_neutralized']:
            print(
                f"   Sanitized docProps: custom.xml={meta['custom_stripped']}, "
                f"core.xml={meta['core_neutralized']}"
            )

    return stats


def main():
    if len(sys.argv) < 3:
        print("Usage: python scripts/clean_template.py <input.pptx> <output.pptx> [--keep-first-master] [--keep-metadata]")
        print()
        print("Options:")
        print("  --keep-first-master  Remove all but the first slide master (experimental)")
        print("  --keep-metadata      Keep docProps (MIP labels, cp:lastModifiedBy, etc.).")
        print("                       Default: strip PII. Use only for tenant-internal repos.")
        sys.exit(1)
    
    input_path = sys.argv[1]
    output_path = sys.argv[2]
    keep_first_master = '--keep-first-master' in sys.argv
    sanitize_metadata = '--keep-metadata' not in sys.argv
    
    if not Path(input_path).exists():
        print(f"Error: File not found: {input_path}")
        sys.exit(1)
    
    # Create output directory if needed
    output_dir = Path(output_path).parent
    output_dir.mkdir(parents=True, exist_ok=True)
    
    stats = clean_template(input_path, output_path, keep_first_master, sanitize_metadata)
    
    # Exit code: 0 on success
    sys.exit(0)


if __name__ == "__main__":
    main()
