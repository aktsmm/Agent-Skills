"""Compare two Azure Update decks semantically, never byte-for-byte."""

from __future__ import annotations

import argparse
import json
import re
import sys
import zipfile
from pathlib import Path

from pptx import Presentation

P14_NS = "http://schemas.microsoft.com/office/powerpoint/2010/main"


def hyperlinks(slide) -> list[str]:
    values = []
    for shape in slide.shapes:
        try:
            address = shape.click_action.hyperlink.address
            if address:
                values.append(address)
        except Exception:
            pass
        if not getattr(shape, "has_text_frame", False):
            continue
        for paragraph in shape.text_frame.paragraphs:
            values.extend(run.hyperlink.address for run in paragraph.runs if run.hyperlink.address)
    return sorted(set(values))


def title(slide) -> str:
    if slide.shapes.title and slide.shapes.title.text.strip():
        return slide.shapes.title.text.strip()
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.text.strip():
            return shape.text.strip().splitlines()[0]
    return ""


def normalized_title(value: str) -> str:
    value = re.sub(r"^【[^】]+】\s*", "", value).strip()
    patterns = (
        r"^(?:Generally Available|General Availability|GA)\s*:\s*",
        r"^(?:Public Preview|Private Preview|Preview)\s*:\s*",
        r"^(?:Retirement|Deprecated|End of Support)\s*:\s*",
        r"^(?:Announcing|Announcement)\s*:\s*",
    )
    for pattern in patterns:
        value = re.sub(pattern, "", value, flags=re.IGNORECASE)
    value = re.sub(r"\s+", " ", value).strip()
    if re.search(r"UPDATE\s*Points|UPDATE.*ポイント", value, re.IGNORECASE):
        return "UPDATE Points"
    return value


def section_snapshot(path: Path) -> list[dict]:
    with zipfile.ZipFile(path) as archive:
        xml = archive.read("ppt/presentation.xml").decode("utf-8")
    sections = []
    for match in re.finditer(r'<p14:section name="([^"]+)"[^>]*>.*?<p14:sldIdLst>(.*?)</p14:sldIdLst>', xml, re.DOTALL):
        sections.append({"name": match.group(1), "slideCount": len(re.findall(r'<p14:sldId id="(\d+)"', match.group(2)))})
    return sections


def snapshot(path: Path) -> dict:
    prs = Presentation(path)
    slides = []
    for index, slide in enumerate(prs.slides, 1):
        texts = [shape.text.strip() for shape in slide.shapes if getattr(shape, "has_text_frame", False) and shape.text.strip()]
        slides.append({
            "index": index,
            "title": normalized_title(title(slide)),
            "layout": slide.slide_layout.name,
            "hidden": slide._element.get("show") == "0",
            "notesPresent": bool(slide.notes_slide.notes_text_frame and slide.notes_slide.notes_text_frame.text.strip()),
            "hyperlinks": hyperlinks(slide),
            "regionTexts": sorted(
                shape.text.strip()
                for shape in slide.shapes
                if shape.name.startswith("RegionStamp") and getattr(shape, "has_text_frame", False) and shape.text.strip()
            ),
        })
    return {"slideCount": len(slides), "slides": slides, "sections": section_snapshot(path)}


def compare_with_accepted(left: dict, right: dict) -> tuple[list[dict], list[dict]]:
    differences = []
    accepted = []
    if left["slideCount"] != right["slideCount"]:
        differences.append({"scope": "deck", "field": "slideCount", "left": left["slideCount"], "right": right["slideCount"]})
    for index in range(min(left["slideCount"], right["slideCount"])):
        left_slide, right_slide = left["slides"][index], right["slides"][index]
        for field in ("title", "layout", "hidden"):
            if field == "title" and left_slide["layout"].startswith("Azure Update Ending") and right_slide["layout"].startswith("Azure Update Ending"):
                if left_slide[field] != right_slide[field]:
                    accepted.append({"scope": f"slide:{index + 1}", "field": field, "reason": "Python writes required formal Ending text", "left": left_slide[field], "right": right_slide[field]})
                continue
            if left_slide[field] != right_slide[field]:
                differences.append({"scope": f"slide:{index + 1}", "field": field, "left": left_slide[field], "right": right_slide[field]})
        missing_links = sorted(set(left_slide["hyperlinks"]) - set(right_slide["hyperlinks"]))
        if missing_links:
            differences.append({"scope": f"slide:{index + 1}", "field": "hyperlinks", "left": left_slide["hyperlinks"], "right": right_slide["hyperlinks"]})
        elif left_slide["hyperlinks"] != right_slide["hyperlinks"]:
            accepted.append({"scope": f"slide:{index + 1}", "field": "hyperlinks", "reason": "Python preserves additional official references", "left": left_slide["hyperlinks"], "right": right_slide["hyperlinks"]})
        if left_slide["regionTexts"] != right_slide["regionTexts"]:
            safe_unknown = left_slide["hidden"] and right_slide["hidden"] and right_slide["regionTexts"] == ["リージョン情報要確認"]
            update_points_continuation = left_slide["title"] == "UPDATE Points" and right_slide["title"] == "UPDATE Points"
            target = accepted if safe_unknown or update_points_continuation else differences
            entry = {"scope": f"slide:{index + 1}", "field": "regionTexts", "left": left_slide["regionTexts"], "right": right_slide["regionTexts"]}
            if safe_unknown:
                entry["reason"] = "Unreviewed hidden Appendix stays neutral instead of claiming Japan unsupported"
            elif update_points_continuation:
                entry["reason"] = "Python avoids the COM defect that stamps an UPDATE Points continuation as region unknown"
            target.append(entry)
        if left_slide["notesPresent"] and not right_slide["notesPresent"]:
            differences.append({"scope": f"slide:{index + 1}", "field": "notesPresent", "left": True, "right": False})
    if left["sections"] != right["sections"]:
        differences.append({"scope": "deck", "field": "sections", "left": left["sections"], "right": right["sections"]})
    return differences, accepted


def compare(left: dict, right: dict) -> list[dict]:
    return compare_with_accepted(left, right)[0]


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--com", required=True, type=Path)
    parser.add_argument("--python", required=True, type=Path)
    parser.add_argument("--result", required=True, type=Path)
    args = parser.parse_args()
    left = snapshot(args.com)
    right = snapshot(args.python)
    differences, accepted = compare_with_accepted(left, right)
    result = {"schemaVersion": 1, "passed": not differences, "com": str(args.com), "python": str(args.python), "differences": differences, "acceptedDifferences": accepted}
    args.result.parent.mkdir(parents=True, exist_ok=True)
    args.result.write_text(json.dumps(result, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    print(json.dumps(result, ensure_ascii=False, indent=2))
    return 0 if not differences else 1


if __name__ == "__main__":
    raise SystemExit(main())
