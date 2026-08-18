"""Build an Azure Update PPTX from reviewed manifests using python-pptx."""

from __future__ import annotations

import argparse
import copy
import fnmatch
import hashlib
import json
import os
import re
import shutil
import sys
import tempfile
import unicodedata
import uuid
import zipfile
from dataclasses import dataclass
from pathlib import Path
from xml.sax.saxutils import escape

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn
from pptx.util import Inches, Pt

ENGINE_VERSION = "1.0.0"
P14_NS = "http://schemas.microsoft.com/office/powerpoint/2010/main"
STATUS_ORDER = {"廃止": 1, "GA": 2, "Preview": 3, "アナウンス": 4, "更新": 5}


@dataclass(frozen=True)
class Paths:
    workspace: Path
    date_folder: Path
    config: Path
    template: Path
    contract: Path
    style: Path
    manifest: Path
    output: Path
    retained: Path
    log: Path


def load_json(path: Path) -> dict:
    with path.open(encoding="utf-8-sig") as handle:
        return json.load(handle)


def sha256(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest().upper()


def normalize_key(value: str) -> str:
    return unicodedata.normalize("NFC", value).casefold()


def rgb(value: str) -> RGBColor:
    match = re.fullmatch(r"#([0-9A-Fa-f]{6})", value)
    if not match:
        raise ValueError(f"Expected #RRGGBB color, got {value!r}")
    raw = match.group(1)
    return RGBColor(int(raw[0:2], 16), int(raw[2:4], 16), int(raw[4:6], 16))


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser()
    parser.add_argument("--workspace-root", required=True, type=Path)
    parser.add_argument("--date-folder", required=True, type=Path)
    parser.add_argument("--output", type=Path)
    parser.add_argument("--contract", type=Path)
    parser.add_argument("--style", type=Path)
    parser.add_argument("--run-id", default=str(uuid.uuid4()))
    parser.add_argument("--result", type=Path)
    return parser.parse_args()


def resolve_paths(args: argparse.Namespace) -> Paths:
    workspace = args.workspace_root.resolve()
    date_folder = args.date_folder if args.date_folder.is_absolute() else workspace / args.date_folder
    date_folder = date_folder.resolve()
    config_path = workspace / ".config" / "config.json"
    config = load_json(config_path)
    template = workspace / config["template"]["folder"] / config["template"]["fileName"]
    python_root = Path(__file__).resolve().parent
    contract = (args.contract or python_root / "template-contract.v1.json").resolve()
    style = (args.style or python_root / "render-style.v1.json").resolve()
    date_token = date_folder.name
    output_name = config["output"]["fileNamePattern"].replace("{year}", str(config["output"]["year"])).replace("{date}", date_token)
    output = (args.output or date_folder / output_name).resolve()
    validation = date_folder / "validation"
    return Paths(
        workspace=workspace,
        date_folder=date_folder,
        config=config_path,
        template=template.resolve(),
        contract=contract,
        style=style,
        manifest=date_folder / "manifest",
        output=output,
        retained=validation / f"{output.stem}-python-openxml.zip",
        log=(getattr(args, "result", None).resolve() if getattr(args, "result", None) else date_folder / "logs" / f"python-build-{args.run_id}.json"),
    )


def slide_texts(slide) -> list[str]:
    return [shape.text.strip() for shape in slide.shapes if getattr(shape, "has_text_frame", False) and shape.text.strip()]


def title_shape(slide):
    if slide.shapes.title is not None:
        return slide.shapes.title
    for shape in slide.shapes:
        if not getattr(shape, "is_placeholder", False):
            continue
        if shape.name.casefold().startswith(("title", "タイトル")):
            return shape
    return next((shape for shape in slide.shapes if getattr(shape, "is_placeholder", False) and getattr(shape, "has_text_frame", False)), None)


def body_placeholder(slide, title_placeholder):
    title_element = title_placeholder._element if title_placeholder is not None else None
    return next(
        (
            shape
            for shape in slide.placeholders
            if shape._element is not title_element and getattr(shape, "has_text_frame", False)
        ),
        None,
    )


def slide_title(slide) -> str:
    shape = title_shape(slide)
    if shape is not None and shape.text.strip():
        return shape.text.strip()
    texts = slide_texts(slide)
    return texts[0] if texts else ""


def matches_role(slide, rule: dict) -> bool:
    layout_name = slide.slide_layout.name
    if rule.get("layoutName") and layout_name != rule["layoutName"]:
        return False
    if rule.get("layoutNameGlob") and not fnmatch.fnmatchcase(layout_name, rule["layoutNameGlob"]):
        return False
    title = slide_title(slide)
    if rule.get("titleRegex") and not re.search(rule["titleRegex"], title, re.IGNORECASE):
        return False
    joined = "\n".join(slide_texts(slide))
    if any(token not in joined for token in rule.get("requiredText", [])):
        return False
    prefixes = rule.get("requiredShapePrefixes", [])
    if prefixes and any(not any(shape.name.startswith(prefix) for shape in slide.shapes) for prefix in prefixes):
        return False
    return True


def validate_template(prs: Presentation, contract: dict) -> dict[str, list[int]]:
    expected = contract["slideSizeEmu"]
    if (prs.slide_width, prs.slide_height) != (expected["width"], expected["height"]):
        raise ValueError("Template slide size does not match template-contract.v1.json")
    if contract.get("notesMasterRequired") and prs.part.package.presentation_part.notes_master is None:
        raise ValueError("Template is missing notesMaster")
    resolved: dict[str, list[int]] = {}
    for role, rule in contract["roles"].items():
        matches = [index for index, slide in enumerate(prs.slides) if matches_role(slide, rule)]
        if role in {"summary", "updatePoints", "body"} and len(matches) != 1:
            raise ValueError(f"Template role {role!r} must resolve exactly once; got {matches}")
        if role in {"cover", "ending"} and not matches:
            raise ValueError(f"Template role {role!r} did not resolve")
        resolved[role] = matches
    return resolved


def profile_value(path: Path, field: str) -> str:
    if not path.exists():
        return ""
    pattern = re.compile(rf"^\|\s*{re.escape(field)}\s*\|\s*(.*?)\s*\|")
    for line in path.read_text(encoding="utf-8").splitlines():
        match = pattern.match(line)
        if match:
            return match.group(1).strip().strip("`")
    return ""


def placeholder_values(paths: Paths, config: dict) -> dict[str, str]:
    profile = paths.workspace / ".config" / "customer-profile.md"
    customer = profile_value(profile, "Customer name")
    system = profile_value(profile, "System name")
    speaker = profile_value(profile, "Speaker")
    customer_base = re.sub(r"\s*(御中|様)$", "", customer).strip()
    date_token = paths.date_folder.name
    display_date = date_token
    if re.fullmatch(r"\d{4}", date_token):
        display_date = f"{config['output']['year']}/{date_token[:2]}/{date_token[2:]}"
    return {
        "{{CUSTOMER}} 御中": f"{customer_base} 御中" if customer_base else "",
        "{{CUSTOMER}}御中": f"{customer_base}御中" if customer_base else "",
        "{{CUSTOMER}}": customer,
        "{{SYSTEM}}向け": f"{system}向け" if system else "",
        "{{SYSTEM}}": system,
        "{{SPEAKER}}": speaker,
        "{{DATE}}": display_date,
    }


def replace_placeholders(prs: Presentation, values: dict[str, str]) -> int:
    count = 0
    ordered = sorted(values.items(), key=lambda item: len(item[0]), reverse=True)
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                original = "".join(run.text for run in paragraph.runs) or paragraph.text
                replaced = original
                for key, value in ordered:
                    if key in replaced:
                        replaced = replaced.replace(key, value)
                        count += 1
                if replaced == original:
                    continue
                if paragraph.runs:
                    paragraph.runs[0].text = replaced
                    for run in paragraph.runs[1:]:
                        run.text = ""
                else:
                    paragraph.text = replaced
    return count


def remove_slide(prs: Presentation, index: int) -> None:
    slide_id = prs.slides._sldIdLst[index]
    prs.part.drop_rel(slide_id.rId)
    prs.slides._sldIdLst.remove(slide_id)


def move_slide(prs: Presentation, current: int, target: int) -> None:
    slide_id = prs.slides._sldIdLst[current]
    prs.slides._sldIdLst.remove(slide_id)
    prs.slides._sldIdLst.insert(target, slide_id)


def hide_slide(slide, hidden: bool = True) -> None:
    if hidden:
        slide._element.set("show", "0")
    elif "show" in slide._element.attrib:
        del slide._element.attrib["show"]


def set_east_asian_font(run, name: str) -> None:
    run.font.name = name
    run._r.get_or_add_rPr().set(qn("a:ea"), name)


def set_text(shape, text: str, size: float | None = None, bold: bool | None = None, font: str = "Meiryo UI") -> None:
    frame = shape.text_frame
    frame.word_wrap = True
    frame.clear()
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    set_east_asian_font(run, font)
    if size is not None:
        run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold


def add_textbox(slide, left, top, width, height, text, *, size=12, bold=False, color="#333333", align=PP_ALIGN.LEFT, font="Meiryo UI"):
    shape = slide.shapes.add_textbox(left, top, width, height)
    frame = shape.text_frame
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    set_east_asian_font(run, font)
    return shape


def add_panel(slide, left, top, width, height, heading: str, body: str, *, fill="#F3F2F1"):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.fill.background()
    set_text(shape, f"{heading}\n{body}", 12, False)
    shape.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(51, 51, 51)
    return shape


def add_link(slide, left, top, width, height, label: str, url: str, name: str):
    shape = add_textbox(slide, left, top, width, height, label, size=8, color="#666666")
    shape.name = name
    run = shape.text_frame.paragraphs[0].runs[0]
    run.hyperlink.address = url
    run.font.underline = True
    shape.click_action.hyperlink.address = url
    return shape


def add_notes(slide, text: str) -> None:
    frame = slide.notes_slide.notes_text_frame
    if frame is not None:
        frame.text = text
        return
    body = parse_xml(
        f"""<p:sp {nsdecls('p', 'a')}><p:nvSpPr><p:cNvPr id=\"2\" name=\"Notes Placeholder 1\"/><p:cNvSpPr txBox=\"1\"/><p:nvPr><p:ph type=\"body\" idx=\"1\"/></p:nvPr></p:nvSpPr><p:spPr/><p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:r><a:t>{escape(text)}</a:t></a:r></a:p></p:txBody></p:sp>"""
    )
    slide.notes_slide._element.cSld.spTree.append(body)


def item_title(item: dict) -> str:
    title = str(item.get("titleJa") or item.get("title") or "").strip()
    patterns = (
        r"^(?:Generally Available|General Availability|GA)\s*:\s*",
        r"^(?:Public Preview|Private Preview|Preview)\s*:\s*",
        r"^(?:Retirement|Deprecated|End of Support)\s*:\s*",
        r"^(?:Announcing|Announcement)\s*:\s*",
        r"^(?:一般提供|一般公開|パブリック\s*プレビュー|プライベート\s*プレビュー|廃止)\s*[:：]\s*",
    )
    for pattern in patterns:
        title = re.sub(pattern, "", title, flags=re.IGNORECASE)
    return title.strip()


def item_label(item: dict) -> str:
    label = str(item.get("label") or "更新").replace("【", "").replace("】", "")
    return label if label in STATUS_ORDER else "更新"


def item_source_url(item: dict) -> str:
    return str(item.get("sourceUrl") or item.get("url") or "")


def resolve_region(item: dict, region_data: dict, *, allow_unknown: bool = False) -> dict:
    direct = str(item.get("japanRegion") or "")
    direct_url = str(item.get("japanRegionUrl") or "")
    if direct:
        return {"status": direct, "source": direct_url}
    regions = region_data.get("regions", region_data.get("services", region_data))
    title = str(item.get("title") or "")
    candidates = [title, item_title(item), str(item.get("id") or "")]
    for key in candidates:
        if key and key in regions:
            info = regions[key]
            return {"status": canonical_region_status(info), "source": str(info.get("source") or "")}
    for info in regions.values() if isinstance(regions, dict) else []:
        if str(item.get("id")) in [str(value) for value in info.get("topicIds", [])]:
            return {"status": canonical_region_status(info), "source": str(info.get("source") or "")}
    if re.search(r"Retirement|廃止|End of Support", title, re.IGNORECASE):
        return {"status": "グローバル", "source": direct_url}
    if allow_unknown:
        return {"status": "リージョン情報要確認", "source": direct_url}
    raise ValueError(f"Reviewed region evidence is missing for {title}")


def canonical_region_status(info: dict) -> str:
    status = str(info.get("status") or info.get("stamp") or "")
    aliases = {
        "both": "Japan East / West 対応",
        "east_only": "Japan East のみ対応",
        "west_only": "Japan West のみ対応",
        "unavailable": "日本リージョン未対応",
        "global": "グローバル",
    }
    if status in aliases:
        return aliases[status]
    if status in aliases.values():
        return status
    if info.get("japanEast") and info.get("japanWest"):
        return "Japan East / West 対応"
    if info.get("japanEast"):
        return "Japan East のみ対応"
    if info.get("japanWest"):
        return "Japan West のみ対応"
    if info.get("global"):
        return "グローバル"
    raise ValueError(f"Unresolved reviewed region status: {info}")


def note_map(notes_data: dict) -> dict[str, str]:
    result: dict[str, str] = {}
    for section in ("weekly", "appendix", "slides"):
        for entry in notes_data.get(section, []):
            title = str(entry.get("title") or "")
            notes = entry.get("notes", entry)
            if isinstance(notes, dict):
                chunks = []
                for key, value in notes.items():
                    if isinstance(value, list):
                        chunks.extend(str(item) for item in value)
                    elif value:
                        chunks.append(str(value))
                text = "\n".join(chunks)
            else:
                text = str(notes)
            if title:
                result[normalize_key(title)] = text
    return result


def build_note(item: dict, notes: dict[str, str], region: dict) -> str:
    title_keys = [normalize_key(str(item.get("title") or "")), normalize_key(item_title(item))]
    supplied = next((notes[key] for key in title_keys if key in notes), "")
    source = item_source_url(item)
    learn = str(item.get("learnUrl") or region.get("source") or "")
    parts = [
        "【概要】", supplied or str(item.get("updateSummary") or item.get("keypoint") or item.get("background") or item_title(item)),
        "【お客様への影響】", str(item.get("customerImpact") or item.get("displayImpact") or "利用状況に応じて影響を確認します。"),
        "【推奨アクション】", str(item.get("action") or item.get("keypoint") or "対象サービスの利用有無と適用条件を確認します。"),
        "【リージョン根拠】", f"{region['status']}: {region.get('source', '')}",
        "【参照】", f"Microsoft Learn 詳細: {learn}", f"Azure Updates 発表: {source}",
        "【想定Q&A】", "適用条件と導入判断は公式ドキュメントの最新情報で確認します。",
    ]
    return "\n".join(part for part in parts if part is not None)


def remove_region_samples(slide) -> None:
    for shape in list(slide.shapes):
        if shape.name.startswith("RegionStamp"):
            slide.shapes._spTree.remove(shape._element)


def add_region_stamp(slide, status: str, style: dict) -> None:
    settings = style["regionStamp"]
    entry = settings["styles"].get(status)
    if not entry:
        raise ValueError(f"Unknown region stamp status: {status}")
    width = Pt(settings["widthPt"])
    height = Pt(settings["heightPt"])
    left = slide.part.package.presentation_part.presentation.slide_width - width - Pt(settings["marginRightPt"])
    top = slide.part.package.presentation_part.presentation.slide_height - height - Pt(settings["marginBottomPt"])
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.name = "RegionStamp_Rendered"
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(entry["background"])
    shape.line.fill.background()
    set_text(shape, entry["text"], settings["fontSizePt"], settings["fontBold"], style["font"]["eastAsian"])
    shape.text_frame.word_wrap = False
    shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.text_frame.paragraphs[0].runs[0].font.color.rgb = rgb(settings["fontColor"])


def populate_body(slide, item: dict, region: dict, style: dict, notes: dict[str, str], contract: dict, *, require_publication_date: bool = True) -> None:
    remove_region_samples(slide)
    title = item_title(item)
    title_limit = int(contract["python"]["maxTitleCharacters"])
    if len(title) > title_limit:
        raise ValueError(f"Title exceeds contract limit ({title_limit}): {title}")
    title_placeholder = title_shape(slide)
    if title_placeholder is not None:
        title_placeholder.left = Inches(0.2)
        title_placeholder.top = Inches(0.15)
        title_placeholder.width = Inches(10.65)
        title_placeholder.height = Inches(0.95)
        set_text(title_placeholder, title, 24, True, style["font"]["eastAsian"])
        title_placeholder.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    label = item_label(item)
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11.2), Inches(0.2), Inches(1.5), Inches(0.42))
    badge.name = "StatusBadge"
    badge.fill.solid()
    badge.fill.fore_color.rgb = rgb(style["statusBadge"].get(label, "#666666"))
    badge.line.fill.background()
    set_text(badge, f"【{label}】", 12, True, style["font"]["eastAsian"])
    badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    badge.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 255, 255)
    body_shape = body_placeholder(slide, title_placeholder)
    body = "\n".join(
        line for line in [
            f"対象：{item.get('targetService') or item.get('category') or 'Azure'}",
            f"更新内容：{item.get('updateSummary') or item.get('keypoint') or item.get('background') or title}",
            f"価値：{item.get('userValue') or item.get('keypoint') or '機能・運用上の選択肢を拡張します。'}",
            f"影響：{item.get('displayImpact') or item.get('customerImpact') or '利用状況に応じて適用条件を確認します。'}",
            f"課金：{item.get('pricing')}" if item.get("pricing") else "",
        ] if line
    )
    if len(body) > 520:
        raise ValueError(f"Body exceeds contract limit (520): {title}")
    if body_shape:
        body_shape.left = Inches(0.6)
        body_shape.top = Inches(1.2)
        body_shape.width = Inches(12.1)
        body_shape.height = Inches(2.75)
        set_text(body_shape, body, 14, False, style["font"]["eastAsian"])
    before = str(item.get("before") or item.get("beforeAfter", {}).get("before") or "従来の構成・制約")
    after = str(item.get("after") or item.get("beforeAfter", {}).get("after") or "更新後の構成・選択肢")
    add_panel(slide, Inches(0.6), Inches(4.25), Inches(5.7), Inches(0.95), "Before", before)
    add_panel(slide, Inches(6.7), Inches(4.25), Inches(5.7), Inches(0.95), "After", after)
    mode = str(item.get("layoutMode") or ("action" if label == "廃止" else "change"))
    mode_heading = {"action": "対応の要点", "technical": "技術の要点", "change": "基礎知識"}.get(mode, "基礎知識")
    mode_body = str(item.get("action") or item.get("keypoint") or item.get("background") or "適用条件と利用シナリオを確認します。")
    add_panel(slide, Inches(0.6), Inches(5.35), Inches(11.8), Inches(0.8), mode_heading, mode_body, fill="#EAF2F8")
    source = item_source_url(item)
    learn = str(item.get("learnUrl") or region.get("source") or "")
    if "learn.microsoft.com" in learn.casefold():
        add_link(slide, Inches(0.6), Inches(6.62), Inches(5.7), Inches(0.2), "参考：Microsoft Learn（詳細）", learn, "MicrosoftLearnReference")
    if source:
        add_link(slide, Inches(6.7), Inches(6.62), Inches(5.7), Inches(0.2), "参考：Azure Updates（発表）", source, "AzureUpdatesReference")
    created = str(item.get("created") or item.get("publishedDate") or "")[:10]
    if require_publication_date and not created:
        raise ValueError(f"Visible publication date is missing: {title}")
    if created:
        add_textbox(slide, Inches(9.0), Inches(6.35), Inches(3.4), Inches(0.18), f"掲載: {created.replace('-', '/')}", size=8, color="#666666", align=PP_ALIGN.RIGHT)
    add_region_stamp(slide, region["status"], style)
    add_notes(slide, build_note(item, notes, region))


def add_summary(prs: Presentation, layout, weekly: list[dict]):
    slide = prs.slides.add_slide(layout)
    title_placeholder = title_shape(slide)
    if title_placeholder is not None:
        set_text(title_placeholder, "Weekly News Topics サマリ", 28, True)
    lines = [f"■ 今週の Weekly New Topics（{len(weekly)}件）"]
    for index, item in enumerate(weekly, 1):
        title = item_title(item)
        short = title if len(title) <= 40 else title[:39] + "…"
        lines.append(f"{index}. 【{item_label(item)}】{short}")
    body = body_placeholder(slide, title_placeholder)
    if body:
        body.left = Inches(0.55)
        body.top = Inches(0.95)
        body.width = Inches(12.0)
        body.height = Inches(5.95)
        summary_size = 10 if len(weekly) > 14 else 12
        set_text(body, "\n".join(lines), summary_size, False)
    add_notes(slide, "\n".join(lines))
    return slide


def chunks(items: list[dict], maximum: int) -> list[list[dict]]:
    if len(items) <= maximum:
        return [items]
    page_count = (len(items) + maximum - 1) // maximum
    base, extra = divmod(len(items), page_count)
    result, offset = [], 0
    for page in range(page_count):
        size = base + (1 if page < extra else 0)
        result.append(items[offset:offset + size])
        offset += size
    return result


def add_update_points(prs: Presentation, layout, pages: list[list[dict]], regions: dict, style: dict, contract: dict):
    headers = contract["updatePoints"]["columns"]
    for page_index, items in enumerate(pages, 1):
        slide = prs.slides.add_slide(layout)
        suffix = f" ({page_index}/{len(pages)})" if len(pages) > 1 else ""
        title_placeholder = title_shape(slide)
        if title_placeholder is not None:
            set_text(title_placeholder, f"今週のUPDATE Points{suffix}", 26, True)
        table = slide.shapes.add_table(len(items) + 1, len(headers), Inches(0.35), Inches(1.25), Inches(12.63), Inches(5.55)).table
        widths = [0.45, 1.65, 3.15, 4.85, 2.53]
        for column, width in zip(table.columns, widths):
            column.width = Inches(width)
        for column_index, header in enumerate(headers):
            cell = table.cell(0, column_index)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0, 120, 212)
        for row_index, item in enumerate(items, 1):
            region = resolve_region(item, regions)
            values = [
                str(row_index + sum(len(page) for page in pages[:page_index - 1])),
                str(item.get("category") or "その他"),
                f"【{item_label(item)}】{item_title(item)}",
                str(item.get("keypoint") or item.get("userValue") or "適用条件を確認"),
                region["status"],
            ]
            for column_index, value in enumerate(values):
                cell = table.cell(row_index, column_index)
                cell.text = value
                cell.text_frame.word_wrap = True
                cell.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                if row_index % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(243, 242, 241)
        add_notes(slide, f"Weekly Topics {len(items)}件のUPDATE Points一覧です。")


def add_appendix_divider(prs: Presentation, layout):
    slide = prs.slides.add_slide(layout)
    title_placeholder = title_shape(slide)
    if title_placeholder is not None:
        set_text(title_placeholder, "Appendix", 30, True)
    add_notes(slide, "参考情報です。通常の説明では非表示にします。")
    hide_slide(slide)
    return slide


def set_ending(slide) -> None:
    text_shapes = [shape for shape in slide.shapes if getattr(shape, "has_text_frame", False)]
    if text_shapes:
        text_shapes[0].name = "Ending-Title"
        set_text(text_shapes[0], "以上", 30, True)
    if len(text_shapes) > 1:
        text_shapes[1].name = "Ending-Subtitle"
        set_text(text_shapes[1], "Azure アップデート情報", 16, False)
    add_notes(slide, "説明を終了します。")


def add_sections(prs: Presentation, sections: list[tuple[str, list]]) -> None:
    presentation = prs.part._element
    ext_list = presentation.find(qn("p:extLst"))
    if ext_list is None:
        ext_list = parse_xml(f"<p:extLst {nsdecls('p')}/>")
        presentation.append(ext_list)
    for existing in list(ext_list):
        if existing.find(f"{{{P14_NS}}}sectionLst") is not None:
            ext_list.remove(existing)
    ext = parse_xml(f'<p:ext {nsdecls("p")} uri="{{521415D9-36F7-43E2-AB2F-B90AF26B5E84}}"/>')
    section_list = parse_xml(f'<p14:sectionLst xmlns:p14="{P14_NS}"/>')
    for name, slides in sections:
        if not slides:
            continue
        section = parse_xml(f'<p14:section xmlns:p14="{P14_NS}" name="{escape(name)}" id="{{{str(uuid.uuid4()).upper()}}}"><p14:sldIdLst/></p14:section>')
        id_list = section.find(f"{{{P14_NS}}}sldIdLst")
        for slide in slides:
            relationship_id = next(
                relationship.rId
                for relationship in prs.part.rels.values()
                if not relationship.is_external and relationship.target_part is slide.part
            )
            source_id = next(value for value in prs.slides._sldIdLst if value.rId == relationship_id)
            node = parse_xml(f'<p14:sldId xmlns:p14="{P14_NS}" id="{source_id.id}"/>')
            id_list.append(node)
        section_list.append(section)
    ext.append(section_list)
    ext_list.append(ext)


def order_items(items: list[dict]) -> list[dict]:
    return sorted(items, key=lambda item: (STATUS_ORDER.get(item_label(item), 9), int(item.get("priority") or 999), item_title(item)))


def build(paths: Paths, run_id: str) -> dict:
    required = [paths.config, paths.template, paths.contract, paths.style, paths.manifest / "fetched-updates.json", paths.manifest / "classification.json", paths.manifest / "region_info_reviewed.json"]
    missing = [str(path) for path in required if not path.exists()]
    if missing:
        raise FileNotFoundError("Missing required inputs: " + ", ".join(missing))
    if not zipfile.is_zipfile(paths.template):
        raise ValueError(f"Template is not an OpenXML ZIP: {paths.template}")
    config = load_json(paths.config)
    contract = load_json(paths.contract)
    style = load_json(paths.style)
    configured_contract = int(config.get("build", {}).get("templateContractVersion", 1))
    configured_style = int(config.get("build", {}).get("renderStyleVersion", 1))
    if configured_contract != int(contract["schemaVersion"]):
        raise ValueError(f"Template contract version mismatch: config={configured_contract}, runtime={contract['schemaVersion']}")
    if configured_style != int(style["schemaVersion"]):
        raise ValueError(f"Render style version mismatch: config={configured_style}, runtime={style['schemaVersion']}")
    classification = load_json(paths.manifest / "classification.json")
    region_data = load_json(paths.manifest / "region_info_reviewed.json")
    notes_data = load_json(paths.manifest / "notes.json") if (paths.manifest / "notes.json").exists() else {}
    weekly = list(classification.get("weekly", []))
    appendix = list(classification.get("appendix", []))
    prs = Presentation(paths.template)
    roles = validate_template(prs, contract)
    replace_placeholders(prs, placeholder_values(paths, config))
    layouts = {layout.name: layout for layout in prs.slide_layouts}
    content_layout = layouts[contract["roles"]["summary"]["layoutName"]]
    body_layout = layouts[contract["roles"]["body"]["layoutName"]]
    cover_slides = [prs.slides[index] for index in roles["cover"]]
    ending_slides = [prs.slides[index] for index in roles["ending"]]
    prototype_slides = [prs.slides[index] for index in roles["summary"] + roles["updatePoints"] + roles["body"]]
    for slide in cover_slides:
        add_notes(slide, "Azure Updateの表紙です。対象期間と説明範囲を確認します。")
    for slide in cover_slides[1:]:
        hide_slide(slide)
    for slide in ending_slides[1:]:
        hide_slide(slide)
    for slide in ending_slides:
        set_ending(slide)
    summary_slide = add_summary(prs, content_layout, weekly)
    weekly_slides = []
    notes = note_map(notes_data)
    for item in weekly:
        slide = prs.slides.add_slide(body_layout)
        populate_body(slide, item, resolve_region(item, region_data), style, notes, contract)
        weekly_slides.append(slide)
    update_slides_before = len(prs.slides)
    add_update_points(prs, content_layout, chunks(weekly, contract["updatePoints"]["maxRowsPerSlide"]), region_data, style, contract)
    update_slides = list(prs.slides)[update_slides_before:]
    appendix_slides = []
    for item in appendix:
        slide = prs.slides.add_slide(body_layout)
        populate_body(slide, item, resolve_region(item, region_data, allow_unknown=True), style, notes, contract, require_publication_date=False)
        hide_slide(slide)
        appendix_slides.append(slide)
    for prototype in prototype_slides:
        remove_slide(prs, list(prs.slides).index(prototype))
    primary_cover = cover_slides[0]
    alternate_covers = cover_slides[1:]
    primary_ending = ending_slides[0]
    alternate_endings = ending_slides[1:]
    desired = [primary_cover, summary_slide, *weekly_slides, *update_slides, *appendix_slides, primary_ending, *alternate_endings, *alternate_covers]
    for target, slide in enumerate(desired):
        current = list(prs.slides).index(slide)
        move_slide(prs, current, target)
    add_sections(prs, [
        ("表紙", [primary_cover]),
        ("サマリ", [summary_slide]),
        ("Weekly New Topics", weekly_slides),
        ("UPDATE Points", update_slides),
        ("Appendix", appendix_slides),
        ("Ending", [primary_ending, *alternate_endings]),
        ("表紙バリエーション", alternate_covers),
    ])
    expected = len(desired)
    if len(prs.slides) != expected:
        raise RuntimeError(f"Slide count mismatch: {len(prs.slides)} != {expected}")
    paths.output.parent.mkdir(parents=True, exist_ok=True)
    paths.retained.parent.mkdir(parents=True, exist_ok=True)
    paths.log.parent.mkdir(parents=True, exist_ok=True)
    temp_dir = Path(tempfile.mkdtemp(prefix="azure-update-python-"))
    temp_output = temp_dir / paths.output.name
    replacement = paths.output.with_suffix(".python-replace.pptx")
    try:
        prs.save(temp_output)
        if not zipfile.is_zipfile(temp_output):
            raise RuntimeError("Python output is not OpenXML")
        with zipfile.ZipFile(temp_output) as archive:
            slide_parts = len([name for name in archive.namelist() if re.fullmatch(r"ppt/slides/slide\d+\.xml", name)])
        if slide_parts != expected:
            raise RuntimeError(f"Saved slide parts mismatch: {slide_parts} != {expected}")
        conflicts = list(paths.output.parent.glob(f"{paths.output.stem} (*){paths.output.suffix}"))
        if conflicts:
            raise RuntimeError("Conflict copies exist: " + ", ".join(path.name for path in conflicts))
        shutil.copy2(temp_output, paths.retained)
        shutil.copy2(temp_output, replacement)
        os.replace(replacement, paths.output)
    finally:
        replacement.unlink(missing_ok=True)
        shutil.rmtree(temp_dir, ignore_errors=True)
    result = {
        "schemaVersion": 1,
        "runId": run_id,
        "state": "built",
        "engine": "python",
        "engineVersion": ENGINE_VERSION,
        "templateContractVersion": contract["schemaVersion"],
        "renderStyleVersion": style["schemaVersion"],
        "output": str(paths.output),
        "outputSha256": sha256(paths.output),
        "retainedOpenXml": str(paths.retained),
        "retainedOpenXmlSha256": sha256(paths.retained),
        "counts": {"weekly": len(weekly), "appendix": len(appendix), "slides": expected},
        "inputHashes": {path.name: sha256(path) for path in required},
    }
    paths.log.write_text(json.dumps(result, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    return result


def main() -> int:
    args = parse_args()
    paths = resolve_paths(args)
    try:
        result = build(paths, args.run_id)
        print(json.dumps(result, ensure_ascii=False, indent=2))
        return 0
    except Exception as error:
        paths.log.parent.mkdir(parents=True, exist_ok=True)
        failure = {"schemaVersion": 1, "runId": args.run_id, "state": "build-failed", "engine": "python", "engineVersion": ENGINE_VERSION, "error": str(error)}
        paths.log.write_text(json.dumps(failure, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
        print(json.dumps(failure, ensure_ascii=False, indent=2), file=sys.stderr)
        return 1


if __name__ == "__main__":
    raise SystemExit(main())
