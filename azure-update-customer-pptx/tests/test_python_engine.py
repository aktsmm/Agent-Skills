from __future__ import annotations

import importlib.util
import json
import shutil
import sys
import tempfile
import unittest
import zipfile
from pathlib import Path

from pptx import Presentation


SKILL_ROOT = Path(__file__).resolve().parents[1]
BUILDER_PATH = SKILL_ROOT / "scripts" / "python" / "build_customer_pptx.py"
SPEC = importlib.util.spec_from_file_location("azure_python_builder", BUILDER_PATH)
BUILDER = importlib.util.module_from_spec(SPEC)
assert SPEC and SPEC.loader
sys.modules[SPEC.name] = BUILDER
SPEC.loader.exec_module(BUILDER)
COMPARE_PATH = SKILL_ROOT / "scripts" / "python" / "compare_pptx.py"
COMPARE_SPEC = importlib.util.spec_from_file_location("azure_python_compare", COMPARE_PATH)
COMPARER = importlib.util.module_from_spec(COMPARE_SPEC)
assert COMPARE_SPEC and COMPARE_SPEC.loader
sys.modules[COMPARE_SPEC.name] = COMPARER
COMPARE_SPEC.loader.exec_module(COMPARER)


class PythonEngineTests(unittest.TestCase):
    def test_balanced_chunks(self):
        pages = BUILDER.chunks([{"id": str(index)} for index in range(11)], 10)
        self.assertEqual([len(page) for page in pages], [6, 5])

    def test_all_region_stamp_statuses(self):
        prs = Presentation(SKILL_ROOT / "assets" / "template" / "azure-update-template.pptx")
        style = BUILDER.load_json(SKILL_ROOT / "assets" / "render-style.v1.json")
        statuses = ["グローバル", "日本リージョン未対応", "Japan East のみ対応", "Japan West のみ対応", "Japan East / West 対応"]
        for status in statuses:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            BUILDER.add_region_stamp(slide, status, style)
            self.assertEqual(slide.shapes[-1].text, style["regionStamp"]["styles"][status]["text"])

    def test_contract_build(self):
        with tempfile.TemporaryDirectory(prefix="azure-python-engine-test-") as temp:
            root = Path(temp)
            (root / ".config").mkdir()
            (root / "template").mkdir()
            (root / "0818" / "manifest").mkdir(parents=True)
            shutil.copy2(SKILL_ROOT / "assets" / "template" / "azure-update-template.pptx", root / "template" / "azure-update-template.pptx")
            config = {
                "build": {"engine": "python", "templateContractVersion": 1, "renderStyleVersion": 1},
                "template": {"folder": "template", "fileName": "azure-update-template.pptx"},
                "output": {"fileNamePattern": "Example_AzureUpdate_{year}{date}.pptx", "year": "2026"},
            }
            (root / ".config" / "config.json").write_text(json.dumps(config, ensure_ascii=False), encoding="utf-8")
            (root / ".config" / "customer-profile.md").write_text(
                "| Customer name | Example Customer |\n| System name | Example System |\n| Speaker | Example Speaker |\n",
                encoding="utf-8",
            )
            items = []
            for index, label in enumerate(("廃止", "GA", "Preview"), 1):
                items.append({
                    "id": str(index),
                    "title": f"Example update {index}",
                    "titleJa": f"サンプル更新 {index}",
                    "label": label,
                    "category": "ネットワーク",
                    "targetService": "Example Service",
                    "updateSummary": f"更新内容 {index}",
                    "userValue": f"価値 {index}",
                    "displayImpact": f"影響 {index}",
                    "before": "従来",
                    "after": "更新後",
                    "keypoint": f"確認ポイント {index}",
                    "sourceUrl": f"https://azure.microsoft.com/updates/?id={index}",
                    "learnUrl": f"https://learn.microsoft.com/azure/example/{index}",
                    "priority": index,
                    "publishedDate": "2026-08-18",
                })
            appendix = [{**items[-1], "id": "4", "title": "Appendix update", "titleJa": "参考更新", "sourceUrl": "https://azure.microsoft.com/updates/?id=4"}]
            manifest = root / "0818" / "manifest"
            (manifest / "fetched-updates.json").write_text(json.dumps({"items": items + appendix}, ensure_ascii=False), encoding="utf-8")
            (manifest / "classification.json").write_text(json.dumps({"weekly": items, "appendix": appendix}, ensure_ascii=False), encoding="utf-8")
            regions = {
                item["title"]: {"status": "Japan East / West 対応", "japanEast": True, "japanWest": True, "verified": True, "source": item["learnUrl"]}
                for item in items
            }
            regions["Appendix update"] = {"status": "グローバル", "global": True, "verified": True, "source": appendix[0]["learnUrl"]}
            (manifest / "region_info_reviewed.json").write_text(json.dumps({"regions": regions}, ensure_ascii=False), encoding="utf-8")
            notes = {"weekly": [{"title": item["title"], "notes": {"summary": f"説明 {item['id']}"}} for item in items], "appendix": []}
            (manifest / "notes.json").write_text(json.dumps(notes, ensure_ascii=False), encoding="utf-8")
            args = type("Args", (), {
                "workspace_root": root,
                "date_folder": Path("0818"),
                "output": None,
                "contract": SKILL_ROOT / "assets" / "template-contract.v1.json",
                "style": SKILL_ROOT / "assets" / "render-style.v1.json",
                "run_id": "fixture-run",
            })()
            result = BUILDER.build(BUILDER.resolve_paths(args), args.run_id)
            output = Path(result["output"])
            self.assertTrue(zipfile.is_zipfile(output))
            prs = Presentation(output)
            self.assertEqual(result["counts"], {"weekly": 3, "appendix": 1, "slides": 12})
            self.assertEqual(len(prs.slides), 12)
            hidden = sum(slide._element.get("show") == "0" for slide in prs.slides)
            self.assertGreaterEqual(hidden, 5)
            note_count = sum(bool(slide.notes_slide.notes_text_frame and slide.notes_slide.notes_text_frame.text.strip()) for slide in prs.slides)
            self.assertEqual(note_count, 12)
            presentation_xml = zipfile.ZipFile(output).read("ppt/presentation.xml").decode("utf-8")
            self.assertIn("Weekly New Topics", presentation_xml)
            self.assertIn("UPDATE Points", presentation_xml)
            all_text = "\n".join(text for slide in prs.slides for text in BUILDER.slide_texts(slide))
            self.assertNotIn("{{CUSTOMER}}", all_text)
            weekly_slide = prs.slides[2]
            link_shapes = [shape for shape in weekly_slide.shapes if shape.name in {"MicrosoftLearnReference", "AzureUpdatesReference"}]
            self.assertEqual(len(link_shapes), 2)
            self.assertTrue(all(shape.click_action.hyperlink.address for shape in link_shapes))
            self.assertEqual(COMPARER.compare(COMPARER.snapshot(output), COMPARER.snapshot(output)), [])


if __name__ == "__main__":
    unittest.main()
