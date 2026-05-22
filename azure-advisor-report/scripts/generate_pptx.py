"""
Azure Advisor Monthly Report - PowerPoint Generator (Template)

Usage:
    py -3 generate_pptx.py --output report.pptx

This script generates a widescreen PowerPoint report with Azure Advisor
recommendations and cost trend analysis.

Customize the DATA section below with actual values from Azure CLI / Cost Management API.
"""
import sys, io, os, argparse
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ============================================================
# CONFIG - Customize these
# ============================================================
FONT_ASCII = 'Calibri'
FONT_EAST_ASIAN = 'BIZ UDPGothic'
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Layout constants (prevent overlap)
HEADER_H = Inches(1.1)       # Header bar height
CONTENT_TOP = Inches(1.3)    # Content starts below header
MARGIN = Inches(0.5)         # Page margin
FULL_W = Inches(12.3)        # Full-width content
LEFT_W = Inches(7.5)         # Left column width (60%)
RIGHT_X = Inches(8.3)        # Right column X position
RIGHT_W = Inches(4.5)        # Right column width (40%)
GAP = Inches(0.3)            # Vertical gap between blocks

# Colors
DARK_BLUE = RGBColor(0x00, 0x32, 0x6E)
MS_BLUE = RGBColor(0x00, 0x78, 0xD4)
RED = RGBColor(0xE7, 0x4C, 0x3C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_BLUE = RGBColor(0xD6, 0xEA, 0xF8)
LIGHT_RED = RGBColor(0xFA, 0xDB, 0xD8)
LIGHT_GREEN = RGBColor(0xD5, 0xF5, 0xE3)

# ============================================================
# HELPER FUNCTIONS
# ============================================================

def _set_font(font, size, bold=False, color=DARK_GRAY):
    """Apply consistent font settings."""
    font.size = Pt(size)
    font.bold = bold
    font.color.rgb = color
    font.name = FONT_ASCII
    font._element.attrib['{http://schemas.openxmlformats.org/drawingml/2006/main}altLang'] = 'ja-JP'
    # Set East Asian font via XML for proper fallback
    from pptx.oxml.ns import qn
    rPr = font._element
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        from lxml import etree
        ea = etree.SubElement(rPr, qn('a:ea'))
    ea.set('typeface', FONT_EAST_ASIAN)

def add_bg(sl, l, t, w, h, c):
    """Add colored rectangle background."""
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background()

def tb(sl, l, t, w, h, txt, fs=12, b=False, c=DARK_GRAY, a=PP_ALIGN.LEFT):
    """Add text box with proper font settings."""
    bx = sl.shapes.add_textbox(l, t, w, h); tf = bx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = txt; _set_font(p.font, fs, b, c); p.alignment = a
    return tf

def ap(tf, txt, fs=12, b=False, c=DARK_GRAY):
    """Add paragraph to text frame."""
    p = tf.add_paragraph(); p.text = txt; _set_font(p.font, fs, b, c)
    p.space_before = Pt(2)

def tbl(sl, l, t, w, h, data, hc=MS_BLUE, col_widths=None):
    """Add table with header coloring and optional column width ratios."""
    r, co = len(data), len(data[0])
    ts = sl.shapes.add_table(r, co, l, t, w, h).table
    # Apply column widths if specified (list of relative proportions)
    if col_widths and len(col_widths) == co:
        total = sum(col_widths)
        for j, cw in enumerate(col_widths):
            ts.columns[j].width = int(w * cw / total)
    for i, rd in enumerate(data):
        for j, ct in enumerate(rd):
            cell = ts.cell(i, j); cell.text = str(ct)
            for p in cell.text_frame.paragraphs:
                _set_font(p.font, 11, bold=(i == 0), color=WHITE if i == 0 else DARK_GRAY)
                if i == 0:
                    p.alignment = PP_ALIGN.CENTER
            if i == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = hc
            elif i % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT_GRAY

def hdr(sl, title, sub=None):
    """Slide header bar."""
    add_bg(sl, Inches(0), Inches(0), SLIDE_W, HEADER_H, DARK_BLUE)
    tb(sl, MARGIN, Inches(0.2), Inches(12), Inches(0.6), title, fs=24, b=True, c=WHITE)
    if sub:
        tb(sl, MARGIN, Inches(0.65), Inches(12), Inches(0.4), sub, fs=13,
           c=RGBColor(0xA0, 0xC0, 0xE0))

def ibox(sl, l, t, w, h, title, lines, bg=LIGHT_BLUE, tc=DARK_BLUE):
    """Insight box with background."""
    add_bg(sl, l, t, w, h, bg)
    tf = tb(sl, l + Inches(0.1), t + Inches(0.05), w - Inches(0.2), h - Inches(0.1),
            title, fs=12, b=True, c=tc)
    for ln in lines:
        ap(tf, ln, fs=11, c=DARK_GRAY)

def add_notes(sl, text):
    """Add speaker notes to a slide."""
    notes_slide = sl.notes_slide
    notes_slide.notes_text_frame.text = text


# ============================================================
# DATA - Replace with actual values
# ============================================================

REPORT_TITLE = "Azure 環境 簡易月次レポート"
REPORT_DATE = "2026年3月31日"
CUSTOMER_NAME = ""  # e.g., "お客様名様"

# Subscription info: [(env_label, sub_name, sub_id_short, annual_usage)]
SUBSCRIPTIONS = [
    ("本番系", "従量課金", "xxxxxxxx-...", "$XXX,XXX"),
    ("評価系", "従量課金（評価）", "yyyyyyyy-...", "$XX,XXX"),
]

# Cost trend: monthly totals (JPY)
COST_MONTHS = ["2025/10", "2025/11", "2025/12", "2026/01", "2026/02", "2026/03"]
COST_PROD = ["0", "0", "0", "0", "0", "0"]           # Production monthly
COST_PROD_MOM = ["-", "+0%", "+0%", "+0%", "+0%", "+0%"]  # MoM change
COST_EVAL = ["0", "0", "0", "0", "0", "0"]           # Evaluation monthly

# Service breakdown: [service, m1, m2, m3, m4, m5, m6, trend]
COST_SERVICES_PROD = [
    ["Service A", "0", "0", "0", "0", "0", "0", "横ばい"],
]

# Advisor recommendations: [recommendation, count, impact, resources]
ADVISOR_COST_PROD = [
    ["VM 適正サイズ化", "N台", "$XX,XXX", "vm-name-1, vm-name-2 ほかN台"],
]
ADVISOR_COST_EVAL = [
    ["VM 適正サイズ化", "N台", "$XX,XXX", "vm-name-1 ほかN台"],
]

ADVISOR_SEC_PROD = [
    ["推奨事項", "N", "High", "resource-1, resource-2 ほかN台"],
]
ADVISOR_SEC_EVAL = [
    ["推奨事項", "N", "High", "resource-1 ほかN台"],
]

ADVISOR_REL_PROD = [
    ["推奨事項", "N", "High", "resource-1 ほかN台"],
]
ADVISOR_REL_EVAL = [
    ["推奨事項", "N", "High", "resource-1 ほかN台"],
]

# Portal links (with tenant ID)
PORTAL_LINKS = [
    ("本番系 Advisor", "https://portal.azure.com/#@{tenantId}/blade/..."),
    ("評価系 Advisor", "https://portal.azure.com/#@{tenantId}/blade/..."),
]


# ============================================================
# SLIDE GENERATION
# ============================================================

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ----- S1: Title -----
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, Inches(0), Inches(0), SLIDE_W, Inches(4), DARK_BLUE)
title_text = f"{CUSTOMER_NAME} {REPORT_TITLE}" if CUSTOMER_NAME else REPORT_TITLE
tb(sl, Inches(1), Inches(0.8), Inches(11), Inches(1.2), title_text, fs=34, b=True, c=WHITE)
tb(sl, Inches(1), Inches(2.0), Inches(11), Inches(0.6),
   "Azure Advisor 推奨事項 / コスト推移分析", fs=16, c=RGBColor(0xA0, 0xC0, 0xE0))
tb(sl, Inches(1), Inches(4.5), Inches(5), Inches(0.4),
   f"レポート作成日: {REPORT_DATE}", fs=14, c=DARK_GRAY)
tb(sl, Inches(1), Inches(5.0), Inches(6), Inches(0.4),
   "データ取得: Azure Advisor / Cost Management API", fs=11, c=DARK_GRAY)

# Subscription table on title slide
sub_data = [["環境", "サブスクリプション名", "Subscription ID", "年間利用額"]]
for env, name, sid, annual in SUBSCRIPTIONS:
    sub_data.append([env, name, sid, annual])
tbl(sl, Inches(6.5), Inches(4.5), Inches(6.3), Inches(1.5), sub_data,
    col_widths=[1, 3, 3, 2])
add_notes(sl, f"表紙スライド。対象: {len(SUBSCRIPTIONS)} サブスクリプション。データ取得日: {REPORT_DATE}")

# ----- S2: Cost Trend -----
sl = prs.slides.add_slide(prs.slide_layouts[6])
hdr(sl, "コスト推移・分析")
tb(sl, MARGIN, CONTENT_TOP, Inches(6), Inches(0.3),
   "本番系 月別コスト推移", fs=13, b=True, c=DARK_BLUE)

trend_data = [[""] + COST_MONTHS, ["合計"] + COST_PROD, ["前月比"] + COST_PROD_MOM]
tbl(sl, MARGIN, CONTENT_TOP + Inches(0.4), FULL_W, Inches(0.8), trend_data)

if COST_SERVICES_PROD:
    svc_top = CONTENT_TOP + Inches(1.4)
    tb(sl, MARGIN, svc_top, Inches(12), Inches(0.3),
       "本番系 サービス別月次推移", fs=12, b=True, c=DARK_BLUE)
    svc_data = [["サービス"] + COST_MONTHS + ["変動"]] + COST_SERVICES_PROD
    svc_h = min(Inches(0.35) * len(svc_data), Inches(2.0))
    tbl(sl, MARGIN, svc_top + Inches(0.4), FULL_W, svc_h, svc_data)

# ----- S3: Cost Optimization -----
sl = prs.slides.add_slide(prs.slide_layouts[6])
hdr(sl, "コスト最適化")

# Production cost table
tb(sl, MARGIN, CONTENT_TOP, Inches(6), Inches(0.3),
   "本番系", fs=13, b=True, c=DARK_BLUE)
cost_prod_data = [["推奨事項", "件数", "年間削減", "対象リソース例"]] + ADVISOR_COST_PROD
cost_h = min(Inches(0.35) * len(cost_prod_data), Inches(2.0))
tbl(sl, MARGIN, CONTENT_TOP + Inches(0.4), FULL_W, cost_h, cost_prod_data,
    col_widths=[5, 1, 2, 5])

# Evaluation cost table
eval_top = CONTENT_TOP + Inches(0.4) + cost_h + GAP
tb(sl, MARGIN, eval_top, Inches(6), Inches(0.3),
   "評価系", fs=13, b=True, c=DARK_BLUE)
cost_eval_data = [["推奨事項", "件数", "年間削減", "対象リソース"]] + ADVISOR_COST_EVAL
eval_h = min(Inches(0.35) * len(cost_eval_data), Inches(1.5))
tbl(sl, MARGIN, eval_top + Inches(0.4), FULL_W, eval_h, cost_eval_data,
    col_widths=[5, 1, 2, 5])
add_notes(sl, "コスト最適化: Advisor Cost カテゴリの推奨事項。年間削減額は Advisor 推定値。")

# ----- S4: Security (Production) -----
sl = prs.slides.add_slide(prs.slide_layouts[6])
hdr(sl, "セキュリティ — 本番系")
sec_data = [["推奨事項", "件数", "影響度", "対象リソース（代表3件 + 残数）"]] + ADVISOR_SEC_PROD
sec_h = min(Inches(0.35) * len(sec_data), Inches(3.5))
tbl(sl, MARGIN, CONTENT_TOP, FULL_W, sec_h, sec_data, hc=RGBColor(0xC0, 0x39, 0x2B),
    col_widths=[5, 1, 1, 6])
add_notes(sl, "セキュリティ（本番系）: Advisor Security カテゴリ。High 影響度を優先対応。")

# ----- S5: Security (Evaluation) -----
sl = prs.slides.add_slide(prs.slide_layouts[6])
hdr(sl, "セキュリティ — 評価系")
sec2_data = [["推奨事項", "件数", "影響度", "対象リソース（代表3件 + 残数）"]] + ADVISOR_SEC_EVAL
sec2_h = min(Inches(0.35) * len(sec2_data), Inches(3.0))
tbl(sl, MARGIN, CONTENT_TOP, FULL_W, sec2_h, sec2_data, hc=RGBColor(0xC0, 0x39, 0x2B),
    col_widths=[5, 1, 1, 6])
add_notes(sl, "セキュリティ（評価系）: 本番と同様に High 影響度を中心に確認。")

# ----- S6: Reliability (Production) -----
sl = prs.slides.add_slide(prs.slide_layouts[6])
hdr(sl, "信頼性・可用性 — 本番系")
rel_data = [["推奨事項", "件数", "影響度", "対象リソース"]] + ADVISOR_REL_PROD
rel_h = min(Inches(0.35) * len(rel_data), Inches(3.5))
tbl(sl, MARGIN, CONTENT_TOP, LEFT_W, rel_h, rel_data, hc=RGBColor(0x29, 0x80, 0xB9),
    col_widths=[5, 1, 1, 5])
ibox(sl, RIGHT_X, CONTENT_TOP, RIGHT_W, Inches(2.5),
     "💡 改善ポイント",
     ["Zone 冗長化やバックアップ設定を", "優先的に確認してください"],
     bg=LIGHT_GREEN, tc=RGBColor(0x1A, 0x5E, 0x3A))
add_notes(sl, "信頼性（本番系）: Advisor HighAvailability カテゴリ。Zone 冗長・Backup を優先。")

# ----- S7: Reliability (Evaluation) -----
sl = prs.slides.add_slide(prs.slide_layouts[6])
hdr(sl, "信頼性・可用性 — 評価系")
rel2_data = [["推奨事項", "件数", "影響度", "対象リソース"]] + ADVISOR_REL_EVAL
rel2_h = min(Inches(0.35) * len(rel2_data), Inches(3.5))
tbl(sl, MARGIN, CONTENT_TOP, LEFT_W, rel2_h, rel2_data, hc=RGBColor(0x29, 0x80, 0xB9),
    col_widths=[5, 1, 1, 5])
add_notes(sl, "信頼性（評価系）: 評価環境の冗長化は必要性を判断の上対応。")

# ----- S8: Appendix -----
sl = prs.slides.add_slide(prs.slide_layouts[6])
hdr(sl, "補足情報")

tf = tb(sl, MARGIN, CONTENT_TOP + Inches(0.2), Inches(6), Inches(4),
        "Azure Portal リンク", fs=14, b=True, c=DARK_BLUE)
for label, url in PORTAL_LINKS:
    ap(tf, "")
    ap(tf, label, fs=12, b=True)
    ap(tf, url, fs=9, c=MS_BLUE)

ibox(sl, Inches(7), CONTENT_TOP + Inches(0.2), Inches(5.8), Inches(4),
     "❗ 免責事項",
     ["本レポートのコスト削減見込額は Azure Advisor の推定値に基づく概算値です。",
      "実際の削減額は利用状況・価格改定・為替変動等により異なる場合があります。",
      "",
      "注意事項:",
      "・推奨事項の件数・対象リソースは日々変動します",
      "・最新の状況は Azure Portal からご確認ください",
      "",
      f"データ取得日: {REPORT_DATE}"],
     bg=LIGHT_RED, tc=RED)
add_notes(sl, "補足情報: Portal 直リンクと免責事項。データ取得日を明示。")


# ============================================================
# SAVE
# ============================================================

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Generate Azure Advisor monthly report PPTX")
    parser.add_argument("--output", "-o", default="report.pptx", help="Output file path")
    args = parser.parse_args()
    out = args.output
    tmp = out + '.tmp.pptx'
    prs.save(tmp)
    if os.path.exists(out):
        os.remove(out)
    os.rename(tmp, out)
    print(f"Saved: {out}")
